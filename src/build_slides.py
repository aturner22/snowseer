"""Build docs/slides.pptx with the Snowseer aesthetic.

Cream background, charcoal text, single rust accent, Inter for headings
and EB Garamond for body text. Mirrors the section structure of
docs/index.html. Two demo slides carry a clearly-labelled dashed
rectangle where the user inserts the corresponding video clip in
Google Slides via Insert -> Video -> Drive.

Run with:
    uv run --with python-pptx python -m src.build_slides
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

ROOT = Path(__file__).resolve().parents[1]
MEDIA = ROOT / "docs/assets/media"
OUT = ROOT / "docs/slides.pptx"

# Identity (matches docs/style/site.css and marp.css).
BG = RGBColor(0xF6, 0xF3, 0xEE)
TEXT = RGBColor(0x1C, 0x1C, 0x1C)
ACCENT = RGBColor(0xB3, 0x4A, 0x25)
MUTE = RGBColor(0x8A, 0x87, 0x80)
HAIRLINE = RGBColor(0xCF, 0xCC, 0xC6)
PANEL = RGBColor(0xEE, 0xEA, 0xE3)

SANS = "Inter"
SERIF = "EB Garamond"
MONO = "JetBrains Mono"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN_X = Inches(0.7)
CONTENT_W = SLIDE_W - 2 * MARGIN_X


# ─── Primitives ────────────────────────────────────────────────────────────


def set_background(slide, color: RGBColor) -> None:
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, *, fill: RGBColor | None = None,
             line: RGBColor | None = None, line_width=Pt(0.75),
             dash: bool = False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = line_width
        if dash:
            ln = shape.line._get_or_add_ln()
            prstDash = etree.SubElement(ln, qn("a:prstDash"))
            prstDash.set("val", "dash")
    return shape


def add_text(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing: float | None = None,
             autosize: bool = False):
    """`runs` is a list of (text, opts) tuples or paragraph specs.

    Paragraph spec: dict with keys {"runs": [...], "align": ..., "space_before": ...}
    Run opts: {"font": str, "size": int, "color": RGBColor, "bold": bool,
              "italic": bool}
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if not autosize:
        tf.auto_size = None

    paragraphs = runs if (runs and isinstance(runs[0], dict)) else [{"runs": runs}]

    for i, pspec in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = pspec.get("align", align)
        if line_spacing is not None:
            p.line_spacing = line_spacing
        if pspec.get("space_before") is not None:
            p.space_before = Pt(pspec["space_before"])
        for r_text, r_opts in pspec["runs"]:
            r = p.add_run()
            r.text = r_text
            r.font.name = r_opts.get("font", SERIF)
            r.font.size = Pt(r_opts.get("size", 18))
            r.font.color.rgb = r_opts.get("color", TEXT)
            r.font.bold = r_opts.get("bold", False)
            r.font.italic = r_opts.get("italic", False)
    return tb


def t(text, **opts):
    """Shorthand: returns (text, opts) for add_text runs."""
    return (text, opts)


def add_chrome(slide, page_num: int, total: int) -> None:
    """Top-left project label, top-right page indicator, bottom rust bar."""
    add_text(
        slide, MARGIN_X, Inches(0.32), Inches(8), Inches(0.3),
        [t("Snowseer  ·  Minimal-Shot Autonomy",
           font=SANS, size=10, color=MUTE)],
    )
    add_text(
        slide, SLIDE_W - MARGIN_X - Inches(2), Inches(0.32),
        Inches(2), Inches(0.3),
        [t(f"{page_num} / {total}", font=SANS, size=10, color=MUTE)],
        align=PP_ALIGN.RIGHT,
    )


def add_title_block(slide, title: str) -> Emu:
    """Slide title in Inter + thin rust accent bar. Returns body-start y."""
    add_text(
        slide, MARGIN_X, Inches(0.85), CONTENT_W, Inches(0.8),
        [t(title, font=SANS, size=32, color=TEXT, bold=True)],
    )
    add_rect(slide, MARGIN_X, Inches(1.6), Inches(0.5), Inches(0.045),
             fill=ACCENT)
    return Inches(1.95)


def add_caption(slide, x, y, w, text: str) -> None:
    add_text(
        slide, x, y, w, Inches(0.6),
        [t(text, font=SANS, size=11, color=MUTE, italic=True)],
        line_spacing=1.3,
    )


# ─── Slide builders ────────────────────────────────────────────────────────


def slide_title(prs, title: str, subtitle: str, footer: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_text(
        slide, Inches(1), Inches(2.7), SLIDE_W - Inches(2), Inches(1.6),
        [t(title, font=SANS, size=84, color=TEXT, bold=True)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    bar_x = int(SLIDE_W // 2) - Inches(0.4)
    add_rect(slide, bar_x, Inches(4.45), Inches(0.8), Inches(0.05), fill=ACCENT)
    add_text(
        slide, Inches(1.5), Inches(4.7), SLIDE_W - Inches(3), Inches(1),
        [t(subtitle, font=SERIF, size=24, color=MUTE, italic=True)],
        align=PP_ALIGN.CENTER, line_spacing=1.4,
    )
    add_text(
        slide, Inches(1), Inches(6.85), SLIDE_W - Inches(2), Inches(0.4),
        [t(footer, font=SANS, size=11, color=MUTE)],
        align=PP_ALIGN.CENTER,
    )
    return slide


def slide_content(prs, title: str, paragraphs, *, page, total, footnote=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_chrome(slide, page, total)
    body_y = add_title_block(slide, title)

    para_specs = []
    for para in paragraphs:
        runs = []
        for chunk in para:
            if isinstance(chunk, str):
                runs.append(t(chunk, font=SERIF, size=18, color=TEXT))
            else:
                text, opts = chunk
                opts = {"font": SERIF, "size": 18, "color": TEXT, **opts}
                runs.append(t(text, **opts))
        para_specs.append({"runs": runs, "space_before": 10})

    add_text(slide, MARGIN_X, body_y, CONTENT_W, Inches(4.5),
             para_specs, line_spacing=1.4)

    if footnote:
        add_text(
            slide, MARGIN_X, Inches(6.55), CONTENT_W, Inches(0.6),
            [t(footnote, font=SANS, size=11, color=MUTE, italic=True)],
            line_spacing=1.4,
        )
    return slide


def slide_image(prs, title: str, image_path: Path, caption: str, *,
                page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_chrome(slide, page, total)
    body_y = add_title_block(slide, title)

    img_max_w = Inches(10.5)
    img_max_h = Inches(3.9)
    pic = slide.shapes.add_picture(
        str(image_path),
        int((SLIDE_W - img_max_w) // 2), body_y + Inches(0.1),
        width=img_max_w,
    )
    if pic.height > img_max_h:
        ratio = img_max_h / pic.height
        pic.height = int(pic.height * ratio)
        pic.width = int(pic.width * ratio)
        pic.left = int((SLIDE_W - pic.width) // 2)

    caption_y = pic.top + pic.height + Inches(0.2)
    caption_h = max(Inches(0.4), Inches(7.1) - caption_y)
    add_text(
        slide, MARGIN_X, caption_y, CONTENT_W, caption_h,
        [t(caption, font=SANS, size=11, color=MUTE, italic=True)],
        line_spacing=1.4, align=PP_ALIGN.CENTER,
    )
    return slide


def slide_demo_placeholder(prs, title: str, video_filename: str,
                           caption: str, *, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_chrome(slide, page, total)
    body_y = add_title_block(slide, title)

    box_w = Inches(10.5)
    box_h = Inches(4.4)
    box_x = int((SLIDE_W - box_w) // 2)
    box_y = body_y + Inches(0.2)
    add_rect(slide, box_x, box_y, box_w, box_h,
             fill=PANEL, line=ACCENT, line_width=Pt(1.5), dash=True)

    add_text(
        slide, box_x, box_y, box_w, box_h,
        [
            {"runs": [t("[ video placeholder ]",
                        font=SANS, size=14, color=ACCENT, bold=True)],
             "align": PP_ALIGN.CENTER},
            {"runs": [t(f"Insert  →  Video  →  Google Drive  →  {video_filename}",
                        font=SANS, size=14, color=MUTE)],
             "align": PP_ALIGN.CENTER, "space_before": 8},
        ],
        anchor=MSO_ANCHOR.MIDDLE,
    )

    caption_y = box_y + box_h + Inches(0.2)
    add_text(
        slide, MARGIN_X, caption_y, CONTENT_W, Inches(0.9),
        [t(caption, font=SANS, size=11, color=MUTE, italic=True)],
        line_spacing=1.4, align=PP_ALIGN.CENTER,
    )
    return slide


def slide_table(prs, title: str, headers, rows, *, page, total, footnote=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_chrome(slide, page, total)
    body_y = add_title_block(slide, title)

    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl_w = CONTENT_W
    tbl_h = Inches(0.55) * n_rows

    table_shape = slide.shapes.add_table(
        n_rows, n_cols, MARGIN_X, body_y + Inches(0.1), tbl_w, tbl_h,
    )
    tbl = table_shape.table

    col_widths = [2.0, 4.5, 2.6, 2.6][:n_cols]
    total_share = sum(col_widths)
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = int(int(tbl_w) * w / total_share)

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG
        _style_cell(cell, h, font=SANS, size=12, color=TEXT, bold=True,
                    border_bottom=TEXT, border_top=None, border_other=None)

    for r_idx, row in enumerate(rows, start=1):
        for c_idx, raw in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG
            chunks = raw if isinstance(raw, list) else [(raw, {})]
            _style_cell(
                cell, chunks,
                font=SERIF, size=14, color=TEXT,
                border_bottom=HAIRLINE if r_idx < n_rows - 1 else None,
                border_top=None, border_other=None,
            )

    if footnote:
        add_text(
            slide, MARGIN_X, Inches(6.55), CONTENT_W, Inches(0.6),
            [t(footnote, font=SANS, size=11, color=MUTE, italic=True)],
            line_spacing=1.4,
        )
    return slide


def _style_cell(cell, content, *, font, size, color, bold=False,
                border_top=None, border_bottom=None, border_other=None):
    cell.margin_left = Inches(0.12)
    cell.margin_right = Inches(0.12)
    cell.margin_top = Inches(0.08)
    cell.margin_bottom = Inches(0.08)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    tf = cell.text_frame
    tf.word_wrap = True

    chunks = [(content, {})] if isinstance(content, str) else content
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    for i, (txt, opts) in enumerate(chunks):
        run = p.add_run() if (i > 0 or p.runs) else p.add_run()
        run.text = txt
        run.font.name = opts.get("font", font)
        run.font.size = Pt(opts.get("size", size))
        run.font.color.rgb = opts.get("color", color)
        run.font.bold = opts.get("bold", bold)
        run.font.italic = opts.get("italic", False)

    _set_cell_borders(cell, top=border_top, bottom=border_bottom,
                      left=border_other, right=border_other)


def _set_cell_borders(cell, *, top=None, bottom=None, left=None, right=None):
    """Hairlines via direct XML (python-pptx doesn't expose cell borders)."""
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ("lnT", "lnB", "lnL", "lnR"):
        for el in tcPr.findall(qn(f"a:{tag}")):
            tcPr.remove(el)
    edges = {"lnL": left, "lnR": right, "lnT": top, "lnB": bottom}
    for tag, color in edges.items():
        if color is None:
            continue
        ln = etree.SubElement(tcPr, qn(f"a:{tag}"))
        ln.set("w", "6350")
        ln.set("cap", "flat")
        ln.set("cmpd", "sng")
        ln.set("algn", "ctr")
        fill = etree.SubElement(ln, qn("a:solidFill"))
        srgb = etree.SubElement(fill, qn("a:srgbClr"))
        srgb.set("val", "{:02X}{:02X}{:02X}".format(color[0], color[1], color[2]))


def slide_diagram(prs, title: str, *, page, total, footnote=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_chrome(slide, page, total)
    body_y = add_title_block(slide, title)

    diagram = (
        "Snow frame                Clear-prior frame\n"
        "      |                            |\n"
        "      |                            v\n"
        "      |                  Mask2Former  (frozen)\n"
        "      |                            |  road mask in prior space\n"
        "      |                            |\n"
        "      +-->  DISK + LightGlue  <----+    (frozen)\n"
        "                  |\n"
        "                  v  correspondences\n"
        "        USAC-MAGSAC homography  (ground-plane biased)\n"
        "                  |\n"
        "                  v  H\n"
        "        warp prior mask -> snow space\n"
        "                  |\n"
        "                  v\n"
        "        fuse over K=3 priors  +  EMA over time\n"
        "                  |\n"
        "                  v\n"
        "        Road overlay on the snow frame"
    )
    add_rect(slide, MARGIN_X, body_y + Inches(0.1),
             CONTENT_W, Inches(4.4), fill=PANEL)
    add_text(
        slide, MARGIN_X + Inches(0.4), body_y + Inches(0.3),
        CONTENT_W - Inches(0.8), Inches(4.2),
        [t(diagram, font=MONO, size=12, color=TEXT)],
        line_spacing=1.2,
    )
    if footnote:
        add_text(
            slide, MARGIN_X, Inches(6.55), CONTENT_W, Inches(0.6),
            [t(footnote, font=SANS, size=11, color=MUTE, italic=True)],
            line_spacing=1.4,
        )
    return slide


def slide_recipe(prs, title: str, items, *, page, total, footnote=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_chrome(slide, page, total)
    body_y = add_title_block(slide, title)

    para_specs = []
    for i, item in enumerate(items, start=1):
        runs = [
            t(f"  {i}.   ", font=SANS, size=18, color=ACCENT, bold=True),
        ]
        for chunk in item:
            if isinstance(chunk, str):
                runs.append(t(chunk, font=SERIF, size=18, color=TEXT))
            else:
                txt, opts = chunk
                opts = {"font": SERIF, "size": 18, "color": TEXT, **opts}
                runs.append(t(txt, **opts))
        para_specs.append({"runs": runs, "space_before": 10})

    add_text(slide, MARGIN_X, body_y + Inches(0.1),
             CONTENT_W, Inches(4.5),
             para_specs, line_spacing=1.4)

    if footnote:
        add_text(
            slide, MARGIN_X, Inches(6.55), CONTENT_W, Inches(0.6),
            [t(footnote, font=SANS, size=11, color=MUTE, italic=True)],
            line_spacing=1.4,
        )
    return slide


def slide_reproduce(prs, *, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_chrome(slide, page, total)
    body_y = add_title_block(slide, "Reproduce")

    code = (
        "git clone https://github.com/aturner22/snowseer\n"
        "cd snowseer\n"
        "uv sync --python 3.12\n"
        "make reproduce"
    )
    code_h = Inches(1.7)
    add_rect(slide, MARGIN_X, body_y + Inches(0.1),
             CONTENT_W, code_h, fill=PANEL)
    add_rect(slide, MARGIN_X, body_y + Inches(0.1),
             Inches(0.04), code_h, fill=ACCENT)
    add_text(
        slide, MARGIN_X + Inches(0.4), body_y + Inches(0.3),
        CONTENT_W - Inches(0.8), code_h,
        [t(code, font=MONO, size=14, color=TEXT)],
        line_spacing=1.4,
    )

    after_y = body_y + Inches(0.1) + code_h + Inches(0.5)
    bullets = [
        ("make track TRACK=<id>", "full pipeline on any registered track."),
        ("make stills", "static-stills precursor (needs MAPILLARY_TOKEN)."),
        ("make pdfs", "render slides.pdf."),
    ]
    para_specs = []
    for cmd, desc in bullets:
        para_specs.append({
            "runs": [
                t(cmd, font=MONO, size=14, color=TEXT, bold=True),
                t("    " + desc, font=SERIF, size=18, color=TEXT),
            ],
            "space_before": 8,
        })
    add_text(slide, MARGIN_X, after_y, CONTENT_W, Inches(2),
             para_specs, line_spacing=1.4)

    add_text(
        slide, MARGIN_X, Inches(6.55), CONTENT_W, Inches(0.5),
        [
            t("Site: ", font=SANS, size=12, color=MUTE),
            t("aturner22.github.io/snowseer", font=MONO, size=12, color=ACCENT),
        ],
    )
    return slide


# ─── Deck ──────────────────────────────────────────────────────────────────


def build() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Build sequence; total page count is set after all slides exist.
    builders: list = []

    builders.append(("title",
        "Snowseer",
        "Achieving minimal-shot autonomy by recognising constants across environments.",
        "SoTA Commission I  ·  May 2026"))

    builders.append(("content",
        "Minimal-shot autonomy",
        [
            ["Minimal-shot autonomy concerns how a system survives in ",
             ("unfamiliar environments", {"bold": True}), "."],
            ["The commonly accepted response is to collect more data and retrain. It does not scale across the long tail of conditions a deployed vehicle, robot, or drone meets in the real world. Snow, dust, smoke, washouts and variable human infrastructure are obvious examples."],
        ],
        "A perception system that depends on having been trained on each new condition will lag every condition it has not yet encountered."))

    builders.append(("content",
        "A snow plough's job is simple",
        [
            [("Sweep the road clear.", {"bold": True})],
            ["The catch is that while the plough is doing it, the road is necessarily invisible. Curbs are buried, lane markings are gone, the boundary between tarmac and verge is no longer defined."],
            ["A self-driving stack trained on traditional road conditions, applied directly to the plough's camera, will report with calibrated confidence that the entire scene is road and should be cleared."],
        ],
        None))

    builders.append(("content",
        "A second move",
        [
            ["For almost every operating environment where autonomy fails for lack of data, an adjacent regime exists, temporally or seasonally or geographically, where data is plentiful and rich, and whose key components remain the same across environments."],
            ["The road that needs to be ploughed this winter is the same road it was in the summer."],
            ["Its appearance has changed. Its position in space and relative to local landmarks has not."],
        ],
        "Snowseer is one demonstration of leveraging structural constants across regimes to achieve minimal-shot autonomy."))

    builders.append(("content",
        "The constants-bridge",
        [
            ["A composition that takes a model trained on regime A, an inference target in regime B, and a known invariant linking the two, and uses the invariant to transfer the model into regime B ",
             ("without retraining", {"bold": True}), "."],
            ["The invariant in this work is geometric. The road sits where it sat last summer, in the same place relative to other landmarks. The shape is general: anatomy across imaging conditions, terrain across illumination, scene structure across weather."],
        ],
        None))

    builders.append(("image",
        "Constants in this work, made visible",
        MEDIA / "nordic_gallivare_matches.png",
        "Snow query (left) and the paired summer prior (right). Green correspondences mark the features that survive the season: gateposts, fence wires, masonry corners, distant roof edges. None land on the road surface itself. The homography fitted to those carries the prior's road mask into the snow frame."))

    builders.append(("recipe",
        "How the system works (per snow frame)",
        [
            [("Pull the live snowy frame", {"bold": True}), " from the plough's camera."],
            [("Pull a clear-season prior", {"bold": True}), " of approximately the same coordinates."],
            [("Match", {"bold": True}), " the two with DISK + LightGlue."],
            ["Estimate a ", ("homography", {"bold": True}), " via USAC-MAGSAC RANSAC."],
            ["Run a Mask2Former segmenter on the ", ("clear prior only", {"italic": True}), "."],
            [("Warp", {"bold": True}), " the road mask into the snow frame and overlay."],
        ],
        "The plough now knows where the road is. No model in the pipeline has been trained on snow."))

    builders.append(("table",
        "Components",
        ["Component", "Role", "Model", "Dataset"],
        [
            ["Feature detector", "Locate keypoints in each image",
             [("DISK ", {"bold": True}), ("(NeurIPS '20)", {"italic": True})],
             "MegaDepth"],
            ["Feature matcher", "Pair keypoints across snow / summer",
             [("LightGlue ", {"bold": True}), ("(ICCV '23)", {"italic": True})],
             "MegaDepth"],
            ["Homography fit", "Robust geometric registration",
             [("USAC-MAGSAC ", {"bold": True}), ("(CVPR '20)", {"italic": True})],
             "n/a"],
            ["Road segmenter", "Road mask on the summer image",
             [("Mask2Former ", {"bold": True}), ("(CVPR '22)", {"italic": True})],
             "Cityscapes"],
            ["Driving captures", "Paired snow + summer Toronto traversals", "n/a",
             [("Boreas ", {"bold": True}), ("(IJRR '23, CC BY 4.0)", {"italic": True})]],
        ],
        "Every learned component is frozen. Snow appears only at inference, as the runtime input."))

    builders.append(("diagram", "The dataflow",
        "A track loader, a prior pool of K=3 nearest summer captures by UTM, and an EMA on the binary mask (alpha = 0.4) wrap the per-still pipeline for the video case."))

    builders.append(("image",
        "One frame, end-to-end",
        MEDIA / "toronto_2021_quad_t005.jpg",
        "Top-left: snow input. Top-right: naive Cityscapes baseline applied directly to snow, painting the road class across snow and sky. Bottom-left: paired summer prior with successful road segmentation. Bottom-right: cross-season overlay produced by warping the prior's road mask through the homography into the snow frame."))

    builders.append(("demo",
        "Demo: Toronto, January 2021",
        "toronto_2021_quad.mp4",
        "14 s drive of a snow-buried residential street. Heavy snow, mid-morning."))

    builders.append(("demo",
        "Demo: Toronto, February 2025",
        "toronto_2025_quad.mp4",
        "34 s drive in active snowfall, late afternoon. Same code, different snow, different drive."))

    builders.append(("image",
        "Nordic stills precursor",
        MEDIA / "nordic_3up.jpg",
        "Three representative pairs from an 18-pair static-stills sweep. Left to right: Gällivare, Kiruna, Luleå. Cross-season overlays in green. Distinct snow scenes, road layouts, lighting and environments, all on a pipeline whose components have never seen snow."))

    builders.append(("content",
        "Limitations",
        [
            [("Some artefacts in the overlay are inherited from the summer prior. ", {"bold": True}),
             "Where the front of the summer capture vehicle is visible, the warped road mask begins a short distance ahead of the snow camera. Where a parked car or other obstacle sits on the road in the prior, the segmenter routes the road class around it and the overlay carries the cutout forward."],
            [("The pipeline is not, currently, real-time. ", {"bold": True}),
             "The matching pass dominates per-frame compute, around 16 s per frame on Mac CPU. Real-time operation needs a substantially faster matcher and segmenter."],
            [("The system is not, currently, deployable arbitrarily. ", {"bold": True}),
             "The current code is geared toward the specific Toronto and nordic demo material. Generalising to any road with Google Street View or a comparable source is feasible (the pipeline is substrate-agnostic in principle), but is a future integration step."],
        ],
        None))

    builders.append(("recipe",
        "Next steps",
        [
            [("Real-time matcher.", {"bold": True}), " Bring per-frame matching from around 16 s on Mac CPU to under 1 s on a deployment-class device."],
            [("Visual place recognition front-end.", {"bold": True}), " Replace GPS-pose lookup with a learned recognition step so the appliance works in GPS-denied environments."],
            [("Multi-source clear-season image bank.", {"bold": True}), " Integrate Mapillary global, Street View, and operator captures so any covered road can be a deployment target."],
            [("Hardware prototype.", {"bold": True}), " A battery-powered processing unit running the live appliance with a simple HUD-style output."],
        ],
        "The snow plough's road-position channel is one consumer of this appliance. The same recipe could power fog, dust, smoke, heavy rain and night driving."))

    builders.append(("reproduce",))

    builders.append(("title",
        "Constants as the bridge",
        "Find what stays the same. Walk across.",
        "Snowseer  ·  SoTA Commission I  ·  May 2026"))

    total = len(builders)
    for i, spec in enumerate(builders, start=1):
        kind = spec[0]
        if kind == "title":
            slide_title(prs, *spec[1:])
        elif kind == "content":
            title, paragraphs, footnote = spec[1], spec[2], spec[3]
            slide_content(prs, title, paragraphs, page=i, total=total,
                          footnote=footnote)
        elif kind == "image":
            title, path, caption = spec[1], spec[2], spec[3]
            slide_image(prs, title, path, caption, page=i, total=total)
        elif kind == "demo":
            title, fname, caption = spec[1], spec[2], spec[3]
            slide_demo_placeholder(prs, title, fname, caption,
                                   page=i, total=total)
        elif kind == "table":
            title, headers, rows, footnote = (spec[1], spec[2], spec[3],
                                              spec[4])
            slide_table(prs, title, headers, rows, page=i, total=total,
                        footnote=footnote)
        elif kind == "diagram":
            slide_diagram(prs, spec[1], page=i, total=total,
                          footnote=spec[2])
        elif kind == "recipe":
            title, items, footnote = spec[1], spec[2], spec[3]
            slide_recipe(prs, title, items, page=i, total=total,
                         footnote=footnote)
        elif kind == "reproduce":
            slide_reproduce(prs, page=i, total=total)
        else:
            raise ValueError(kind)

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    out = build()
    size_mb = out.stat().st_size / (1024 * 1024)
    print(f"wrote {out.relative_to(ROOT)}  ({size_mb:.1f} MB)")
